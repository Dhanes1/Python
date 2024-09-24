from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Mastering PowerPoint Transitions"
subtitle.text = "A guide to enhancing your presentations"

# Slide 2: Introduction
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction"
content.text = "Brief introduction to the importance of transitions in presentations. Mention the impact on audience engagement and flow."

# Add more slides following a similar pattern...

# Save the presentation
presentation.save("PowerPoint_Transitions_Guide.pptx")
